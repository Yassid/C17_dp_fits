#!/usr/bin/env python3
"""Editable PowerPoint version of the analysis-overview deck (python-pptx).

Mirrors make_deck.py; embeds the same PNGs from plots/.  Output:
    Presentations/C17_analysis_overview.pptx
"""
import os
import csv

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN

REPO = os.path.join(os.path.dirname(__file__), "..")
PLOTS = os.path.join(REPO, "plots")
RESULTS = os.path.join(REPO, "results")
OUT = os.path.join(REPO, "Presentations", "C17_analysis_overview.pptx")

MAROON = RGBColor(0x6F, 0x1D, 0x46)
ACCENT = RGBColor(0xB3, 0x47, 0x7D)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
INK = RGBColor(0x1A, 0x1A, 0x1A)
GREY = RGBColor(0x55, 0x55, 0x55)
PINK = RGBColor(0xE8, 0xC9, 0xDA)

SW, SH = 13.333, 7.5

LSTATE_LABEL = {
    "ex2.763_1half-": "2.763 1/2-", "ex2.980_3half+": "2.980 3/2+",
    "ex3.661_NEW": "3.661 new", "ex4.231_3half+": "4.231 3/2+",
    "ex4.841_1half+": "4.841 1/2+", "ex5.91_3half+": "5.91 3/2+",
    "ex6.30_5half+": "6.30 5/2+",
}


def png_size(p):
    with open(p, "rb") as f:
        b = f.read(24)
    return int.from_bytes(b[16:20], "big"), int.from_bytes(b[20:24], "big")


def blank(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])


def header(slide, title, kicker=None):
    r = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0),
                               Inches(SW), Inches(0.95))
    r.fill.solid(); r.fill.fore_color.rgb = MAROON; r.line.fill.background()
    r.shadow.inherit = False
    a = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0.95),
                               Inches(SW), Inches(0.05))
    a.fill.solid(); a.fill.fore_color.rgb = ACCENT; a.line.fill.background()
    a.shadow.inherit = False
    tb = slide.shapes.add_textbox(Inches(0.3), Inches(0.13), Inches(9.8),
                                  Inches(0.7))
    p = tb.text_frame.paragraphs[0]
    run = p.add_run(); run.text = title
    run.font.size = Pt(27); run.font.bold = True; run.font.color.rgb = WHITE
    if kicker:
        kb = slide.shapes.add_textbox(Inches(9.6), Inches(0.28), Inches(3.4),
                                      Inches(0.5))
        kp = kb.text_frame.paragraphs[0]; kp.alignment = PP_ALIGN.RIGHT
        run = kp.add_run(); run.text = kicker
        run.font.size = Pt(13); run.font.italic = True; run.font.color.rgb = WHITE
    fb = slide.shapes.add_textbox(Inches(0.3), Inches(7.05), Inches(9.5),
                                  Inches(0.35))
    run = fb.text_frame.paragraphs[0].add_run()
    run.text = "¹⁶C(d,p)¹⁷C  ·  FRIB a1975  ·  Y. Ayyad (USC)"
    run.font.size = Pt(9); run.font.color.rgb = GREY


def text(slide, x, y, w, h, items, size=16, bold=False, color=INK,
         mono=False, italic=False, align=PP_ALIGN.LEFT, after=6):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame; tf.word_wrap = True
    first = True
    for it in items:
        lvl = 0; s = it
        while s.startswith(">"):
            lvl += 1; s = s[1:].lstrip()
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.alignment = align
        prefix = "" if mono else ("•  " if lvl == 0 else "–  ")
        run = p.add_run(); run.text = prefix + s
        f = run.font
        f.size = Pt(size - 2 * lvl); f.bold = bold; f.italic = italic
        f.color.rgb = color
        if mono:
            f.name = "Courier New"
        p.space_after = Pt(after)
        if lvl:
            p.level = lvl
    return tb


def pic(slide, path, bx, by, bw, bh):
    iw, ih = png_size(path); ar = iw / ih
    if bw / bh > ar:
        h = bh; w = bh * ar; x = bx + (bw - w) / 2; y = by
    else:
        w = bw; h = bw / ar; x = bx; y = by + (bh - h) / 2
    slide.shapes.add_picture(path, Inches(x), Inches(y), Inches(w), Inches(h))


def read_lscan():
    return list(csv.DictReader(open(os.path.join(RESULTS, "fine", "L_scan.csv"))))


def build():
    prs = Presentation()
    prs.slide_width = Inches(SW); prs.slide_height = Inches(SH)
    L = read_lscan()

    # ---- 1. title -----------------------------------------------------
    s = blank(prs)
    bg = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0),
                            Inches(SW), Inches(SH))
    bg.fill.solid(); bg.fill.fore_color.rgb = MAROON; bg.line.fill.background()
    bg.shadow.inherit = False
    text(s, 0.5, 1.7, 12.3, 1.2,
         ["Spectroscopy of unbound ¹⁷C states"],
         size=40, bold=True, color=WHITE, align=PP_ALIGN.CENTER, after=0)
    text(s, 0.5, 2.9, 12.3, 0.8,
         ["via ¹⁶C(d,p)¹⁷C at FRIB (a1975)"],
         size=26, color=WHITE, align=PP_ALIGN.CENTER, after=0)
    text(s, 0.5, 4.0, 12.3, 0.6,
         ["Angular distributions · absolute normalization · L assignments"],
         size=17, color=WHITE, italic=True, align=PP_ALIGN.CENTER, after=0)
    text(s, 0.5, 4.9, 12.3, 0.5,
         ["Y. Ayyad   —   IGFAE, Universidade de Santiago de Compostela"],
         size=16, color=WHITE, align=PP_ALIGN.CENTER, after=0)
    text(s, 0.5, 5.5, 12.3, 0.4, ["30 May 2026"], size=14, color=WHITE,
         align=PP_ALIGN.CENTER, after=0)
    text(s, 0.5, 6.4, 12.3, 0.4,
         ["spectroscopic reference: Movilla & Ayyad (2026-02)"],
         size=11, color=PINK, align=PP_ALIGN.CENTER, after=0)

    # ---- 2. experiment ------------------------------------------------
    s = blank(prs); header(s, "The experiment", "reaction & setup")
    text(s, 0.5, 1.3, 7.6, 5.4, [
        "¹⁶C beam in the AT-TPC, D₂ active target (FRIB a1975).",
        "188.33 MeV total = 11.77 MeV/u;  Q(g.s.) = −1.498 MeV.",
        "Transfer populates the bound region + seven unbound resonances.",
        "Goal: per-state dσ/dΩ on an absolute scale",
        ">→ transferred-L from FRESCO ADWA shape comparison.",
        "Data reduced with ATTPCROOT (InterpSolver); cont. of C16dp_fits.",
    ], size=16, after=10)
    text(s, 9.0, 1.3, 3.9, 0.4, ["Unbound resonances fitted"], size=14,
         bold=True, color=MAROON)
    text(s, 9.0, 1.9, 4.0, 4.5, [
        "Ex(MeV)  Jpi    G(keV)", "----------------------",
        "2.763    1/2-      50", "2.980    3/2+     200",
        "3.661    new      398", "4.231    3/2+     500",
        "4.841    1/2+     300", "5.91     3/2+    1500",
        "6.30     5/2+     396",
    ], size=13, mono=True, after=2)

    # ---- 3. pipeline --------------------------------------------------
    s = blank(prs); header(s, "Analysis pipeline", "method")
    text(s, 0.5, 1.3, 12.3, 5.4, [
        "1. Spectral fit (C16_pd_AngBins.C): 12 × 2.5° bins, 10–40°.",
        ">3 Gaussians (bound) + 7 penetrability-modified Breit–Wigners",
        ">+ 1n/2n phase space + linear bg;  analytic Voigt lineshape.",
        ">Inclusive fit pins positions/widths; per-bin frees amplitudes.",
        "2. Efficiency: 8 Ex-tables, bilinear (Ex,θ); floor + gradient cuts.",
        "3. Absolute normalization fixed by the ¹⁶C(d,d) elastic channel.",
        "4. L scan: each AD vs FRESCO L=0–3 ADWA shapes; best L by χ²/ν.",
        "5. Systematics: 82-variant campaign → per-state envelope.",
    ], size=16, after=10)

    # ---- 4. inclusive -------------------------------------------------
    s = blank(prs); header(s, "Inclusive Ex spectrum: the global fit",
                           "spectral fit")
    pic(s, os.path.join(PLOTS, "plots_inclusive.png"), 0.5, 1.2, 7.4, 5.5)
    text(s, 8.3, 1.3, 4.6, 1.6, [
        "A localized +2–4σ excess near 0.4–1.1 MeV is not absorbed "
        "by the nominal model.",
        "Adding a contaminant Gaussian (~0.45 MeV) restores an acceptable fit:",
    ], size=13)
    text(s, 8.3, 3.5, 4.6, 1.4, [
        "cuts     nominal  +contam", "-------------------------",
        "loose     1.78     0.84", "strict    1.88     0.97",
        "         (chi2/nu)",
    ], size=12.5, mono=True, after=2)
    text(s, 8.3, 5.3, 4.6, 1.5, [
        "Survives strict cuts (0.450→0.451 MeV) ⇒ not an AT-TPC edge "
        "artifact.",
        "Carried as a model-dependence systematic.",
    ], size=12.5)

    # ---- 5. bound states ----------------------------------------------
    s = blank(prs); header(s, "Bound states: angular distributions", "comparison")
    pic(s, os.path.join(PLOTS, "bound_states_ad.png"), 0.4, 1.15, 8.0, 5.6)
    text(s, 8.8, 1.3, 4.2, 5.4, [
        "Three bound ¹⁷C states, as a cross-check of the machinery:",
        ">g.s. 3/2⁺",
        ">0.217 MeV 1/2⁺",
        ">0.335 MeV 5/2⁺",
        "Forward-peaked g.s. and 0.217 shapes consistent with the expected "
        "ℓ content.",
        "The 5/2⁺ is only weakly populated — points sit near the "
        "detection floor (large bars), shown for completeness.",
        "Validate efficiency + absolute scale before the unbound region.",
    ], size=13, after=8)

    # ---- 6. normalization idea ----------------------------------------
    s = blank(prs); header(s, "Absolute normalization: the idea", "normalization")
    text(s, 0.5, 1.3, 12.0, 0.6,
         ["Counts → cross section (fixed luminosity constant):"],
         size=16)
    text(s, 0.5, 2.0, 12.3, 0.9,
         ["dσ/dΩ = Y / ( N_beam · N_tgt · ΔΩ ) "
          "× 10  /  ε"],
         size=22, bold=True, align=PP_ALIGN.CENTER, after=0)
    text(s, 0.5, 3.1, 12.3, 3.4, [
        "N_beam = 1.6146×10⁵,  N_tgt = 0.019632,  "
        "ΔΩ = 2π(cosθ_lo − cosθ_hi).",
        "Cross-check: the ¹⁶C(d,d) elastic channel shares the entrance "
        "partition (same beam, D₂ target, luminosity).",
        ">An absolute OM elastic cross section that reproduces the data "
        "validates the constant — no free normalization.",
        "Self-contained partial-wave OM solver (om_elastic.py): Numerov + "
        "Coulomb matching; vs Rutherford to 10⁻³.",
        ">E_lab(d) = 23.55 MeV, E_c.m. = 20.92 MeV.",
    ], size=15, after=10)

    # ---- 7. normalization result --------------------------------------
    s = blank(prs); header(s, "Absolute normalization: result", "normalization")
    pic(s, os.path.join(PLOTS, "om_elastic_dd.png"), 0.4, 1.15, 7.6, 5.6)
    text(s, 8.4, 1.3, 4.6, 0.4, ["Forward-lobe ⟨data/OM⟩:"],
         size=14, bold=True)
    text(s, 8.4, 1.9, 4.6, 1.8, [
        "OMP            chi2/nu  d/OM", "---------------------------",
        "ADWA(transfer)   398   0.864", "DA1p (global)    105   0.886",
        "free 5-par fit    32   0.978",
    ], size=12.5, mono=True, after=2)
    text(s, 8.4, 3.9, 4.6, 2.8, [
        "Absolute scale consistent with unity — no rescale.",
        "BUT the transfer ADWA potential itself sits ~14% low.",
        "Honest figure = cross-potential spread ≈ 0.86–0.98 (~12%), "
        "a GLOBAL scale uncertainty (not in per-point errors).",
    ], size=12.5)

    # ---- 8. error model -----------------------------------------------
    s = blank(prs); header(s, "From counts to dσ/dΩ: error model",
                           "uncertainties")
    text(s, 0.5, 1.5, 12.3, 0.9,
         ["δ(dσ/dΩ) = (dσ/dΩ) · "
          "√[ (δY/Y)² + (δε/ε)² ]"],
         size=21, bold=True, align=PP_ALIGN.CENTER, after=0)
    text(s, 0.5, 2.7, 12.3, 4.0, [
        "δY = Minuit amplitude error (χ² fit to √N-weighted "
        "histograms) — already carries Poisson stats.",
        "FIXED (2026-05-30): a previous √(1/Y + (δY/Y)²) "
        "double-counted Poisson (√2 inflation).",
        ">e.g. g.s. 10–12.5°:  0.433 → 0.356 mb/sr.",
        "Three uncertainty layers, kept separate:",
        ">statistical (above) + systematic (82-variant envelope) "
        "+ global scale (~12% OMP spread).",
    ], size=15.5, after=10)

    # ---- 9. unbound grid ----------------------------------------------
    s = blank(prs); header(s, "Unbound resonances: dσ/dΩ + "
                           "systematic band", "angular distributions")
    pic(s, os.path.join(PLOTS, "unbound_ad_grid.png"), 0.3, 1.15, 12.7, 5.2)
    text(s, 0.5, 6.45, 12.3, 0.5, [
        "Black: data ± corrected statistical error.   Gray: min–max "
        "across the 82 spectral-fit variants.   Fits use θ_c.m. ≤ 30°.",
    ], size=11, color=GREY, align=PP_ALIGN.CENTER)

    # ---- 10. L method -------------------------------------------------
    s = blank(prs); header(s, "L assignment: method", "L values")
    text(s, 0.5, 1.3, 12.3, 5.4, [
        "FRESCO 7-state DWBA with the Johnson–Soper ADWA effective deuteron "
        "OMP (sum of Koning–Delaroche p+¹⁶C and n+¹⁶C at E_d/2).",
        "Four templates L=0,1,2,3 (neutron in 3s₁/₂, 2p₁/₂, "
        "2d₃/₂/2d₅/₂, 1f₅/₂); weakly-bound approx "
        "b_e = 0.5 MeV.",
        "Each AD fit to each L shape (single normalization); best L by "
        "χ²/ν on θ_c.m. ≤ 30°.",
        "SHAPE-ONLY: the WBA reproduces the L-dependent shape but not the "
        "magnitude ⇒ C²S are effective, not publication SFs",
        ">(needs CDCC / pole-residue treatment).",
    ], size=15.5, after=12)

    # ---- 11. FRESCO overlays ------------------------------------------
    s = blank(prs); header(s, "L assignment: FRESCO shape overlays", "L values")
    pic(s, os.path.join(PLOTS, "plots_fresco_unbound.png"), 0.6, 1.15, 12.1, 5.2)
    text(s, 0.5, 6.45, 12.3, 0.5, [
        "Curves = FRESCO ADWA at each state's literature single-particle config "
        "(normalization-fit, full 10–40°). The data-driven all-L scan on "
        "θ_c.m. ≤ 30° is the next slide.",
    ], size=11, color=GREY, align=PP_ALIGN.CENTER)

    # ---- 12. L results ------------------------------------------------
    s = blank(prs); header(s, "L assignment: results", "L values")
    pic(s, os.path.join(PLOTS, "lscan_bars.png"), 0.3, 1.3, 7.4, 5.0)
    text(s, 7.9, 1.3, 5.1, 0.4, ["best L (corrected errors):"], size=13,
         bold=True)
    text(s, 7.9, 1.9, 5.1, 3.4, [
        "Ex      best L      next L", "-----------------------------",
        "2.763   L=3 (1.41)  L=2 (1.42)", "2.980   L=3 (2.73)  L=2 (3.05)",
        "3.661  *L=2 (1.38)  L=3 (1.45)", "4.231  *L=2 (0.58)  L=3 (1.38)",
        "4.841  *L=0 (3.09)  L=1 (7.35)", "5.91    L=2 (3.55)  L=0 (3.73)",
        "6.30    L=0 (0.34)  L=3 (0.40)",
    ], size=12, mono=True, after=2)
    text(s, 7.9, 5.1, 5.1, 1.6, [
        "χ²/ν refreshed after the error fix.",
        "All best-L picks & tie flags UNCHANGED;",
        ">χ²/ν rose ×1.3–1.5 (every L scales together).",
    ], size=12)

    # ---- 13. state-by-state -------------------------------------------
    s = blank(prs); header(s, "L assignment: state-by-state", "L values")
    text(s, 0.5, 1.3, 12.3, 5.4, [
        "4.231 3/2⁺ / 4.841 1/2⁺: decisive (L=2 d-wave; L=0 s-wave) "
        "— next L worse by factor ≳ 2; match the PDF.",
        "3.661 (new): L=2 over L=3 (1.38 vs 1.45) — not decisive; possible "
        "sd–pf intruder.",
        "2.763 1/2⁻: three-way near-tie (L=3/2/1); AD too flat to separate.",
        "2.980 3/2⁺: no template fits well; 17.5° AD peak not "
        "reproduced; fit wants E₀ ~120 keV above the SM prior.",
        "6.30 5/2⁺: L=0 has best χ²/ν but is FORBIDDEN "
        "(0⁺⊗s₁/₂ = 1/2⁺); allowed L=2/3 are tied.",
        "5.91 3/2⁺: L=2 marginal over L=0; AD poorly constrained "
        "past 30°.",
    ], size=14.5, after=11)

    # ---- 14. systematics ----------------------------------------------
    s = blank(prs); header(s, "Systematic envelope (82 variants)", "systematics")
    pic(s, os.path.join(PLOTS, "yield_envelope_fine.png"), 0.4, 1.15, 7.6, 5.6)
    text(s, 8.4, 1.3, 4.6, 0.4, ["Integrated-yield envelope:"], size=13,
         bold=True)
    text(s, 8.4, 1.9, 4.6, 3.4, [
        "state        -%     +%", "---------------------",
        "g.s. 3/2+   -23    +0", "2.763       -56   +32",
        "2.980       -28  +179", "3.661       -15   +17",
        "4.231       -14   +48", "4.841       -27   +11",
        "5.91         -8   +22", "6.30         -9    +3",
    ], size=12, mono=True, after=2)
    text(s, 8.4, 5.2, 4.6, 1.5, [
        "Largest model-dependence: 2.763 & 2.980 (contaminant degeneracy).",
        "6.30 robust (<10%).",
    ], size=12)

    # ---- 15. summary --------------------------------------------------
    s = blank(prs); header(s, "Summary")
    text(s, 0.5, 1.3, 12.3, 5.4, [
        "Per-state absolute dσ/dΩ for 3 bound + 7 unbound ¹⁷C "
        "states from ¹⁶C(d,p) at 11.77 MeV/u.",
        "Absolute scale validated vs ¹⁶C(d,d) elastic; normalization "
        "uncertainty = ~12% OMP spread (global).",
        "Error model corrected: per-point stats no longer double-counted; "
        "statistical + systematic + scale layers separated.",
        "L assignments robust to the fix: 4.231 (L=2) & 4.841 (L=0) decisive; "
        "3.661/2.763/5.91 ambiguous; 2.980 & 6.30 problematic.",
        "Deliverable (this repo): yields, dσ/dΩ, systematic envelope. "
        "Final L / C²S done downstream (CDCC-grade).",
    ], size=15.5, after=12)

    # ---- 16. open items -----------------------------------------------
    s = blank(prs); header(s, "Open items")
    text(s, 0.5, 1.3, 12.3, 4.6, [
        "Fold the ~12% elastic-normalization scale into the published "
        "dσ/dΩ / yield bands (common factor).",
        "Continuum-form tests (quadratic/exp bg, PS-shape) to show the "
        "0.45 MeV contaminant is not bg flexibility.",
        "Physically identify the 0.45 MeV excess (vertex_z, θ_lab "
        "fingerprint of a non-D₂ channel?).",
        "CDCC / pole-residue treatment for publication-grade absolute C²S.",
    ], size=16, after=12)
    text(s, 0.5, 6.2, 12.3, 0.6, [
        "Refs: Pereira-Lopez et al., PLB 811 (2020) 135939; Movilla & Ayyad "
        "(2026); Pang et al., arXiv:1606.01507 (DA1p).",
    ], size=10, color=GREY)

    # ---- 17. backup grid ----------------------------------------------
    s = blank(prs); header(s, "Backup: full χ²/ν grid", "L values")
    hdr = "state         L=0     L=1     L=2     L=3    best"
    rows = [hdr, "-" * len(hdr)]
    for r in L:
        lab = LSTATE_LABEL.get(r["state"], r["state"])
        best = int(r["bestL"])
        cells = []
        for i in range(4):
            v = float(r[f"c2v_L{i}"])
            cells.append(("*" if i == best else " ") + f"{v:5.2f}")
        tie = ("  (tie L=%s)" % r["tieL"]) if r["tie"] == "1" else ""
        rows.append(f"{lab:<11s}" + " ".join(cells) + f"   L={best}{tie}")
    text(s, 0.6, 1.4, 12.3, 4.6, rows, size=14, mono=True, after=4)
    text(s, 0.6, 6.0, 12.3, 0.9, [
        "* = best L.  Fits on theta_cm <= 30 deg, corrected stat errors "
        "(2026-05-30).",
        "tie = runner-up within 0.3 in chi2/nu (shapes indistinguishable at "
        "current statistics).",
    ], size=11, color=GREY, after=2)

    # ---- 18. backup: non-resonant background --------------------------
    s = blank(prs)
    header(s, "Backup: non-resonant background treatment", "background")
    pic(s, os.path.join(PLOTS, "fine", "plots_bins.png"), 0.2, 1.15, 8.0, 5.6)
    pic(s, os.path.join(PLOTS, "bg_vs_angle.png"), 8.45, 1.2, 4.7, 3.0)
    text(s, 8.4, 4.45, 4.8, 2.5, [
        "In the fit: 1n PS + linear bg + 2n PS, amplitudes free in EVERY "
        "angular bin (PS templates rebuilt per bin).",
        "From ~30° the transfer peaks fade and the continuum dominates the "
        "shape.",
        "Scaling each term 0.3×–3× moves dσ/dΩ by <2% (bg), ~0% (PS) — flat "
        "in angle, no rise past 30°.",
        "Also in the 82-variant envelope; L-scan uses θ ≤ 30°. Back-angle "
        "limit is stats/efficiency, not bg.",
    ], size=11)

    prs.save(OUT)
    print("wrote", OUT, "(", len(prs.slides), "slides )")


if __name__ == "__main__":
    build()
